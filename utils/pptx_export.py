"""utils/pptx_export.py — slide decks on the Tabula magna 2026 template.

Three decks, one discipline: **use the template, add nothing to it.**

* Every deck opens with the *Tabula prima* title layout.
* Structure (titles, section breaks, captions, footers) goes into the
  layouts' own PLACEHOLDERS — *Sectio plena I–IV* for the numbered sections,
  *Tabella risultati* for every list (its native TABLE placeholder),
  *Due contenuti* for the per-person pages, *Messaggio chiave* for the PhD
  clock. No hand-drawn header bars, tiles or footer boxes: formatting comes
  from the template, so when the template evolves the decks follow.
* The ONE exception is the two timeline figures (conference submissions and
  the publication pipeline): those are data graphics with no layout
  equivalent, drawn on *Solo titolo* using only theme colours.

Decks:
* **Meeting** — the weekly delta; whoever is in the room knows the context.
* **Review** — cumulative status per deliverable, for funders/partners.
* **My status** — one person's portfolio for 1:1s and PhD reviews.

Template file: ``assets/maic_template_2026.pptx`` (the v2.0 .potx converted to
.pptx — same package, different declared content type). Palette and fonts are
read from its theme and mirrored in the constants below.
"""

from __future__ import annotations

import datetime
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

TEMPLATE = Path(__file__).resolve().parent.parent / "assets" / "maic_template_2026.pptx"

# ── MAIC LAB 2026 palette (theme1.xml of the template) ────────────────────────
INK     = RGBColor(0x22, 0x28, 0x2E)   # dk1  — body text
MUTED   = RGBColor(0x8D, 0x95, 0x9E)   # accent5 — secondary text
RULE    = RGBColor(0xE9, 0xEB, 0xEE)   # lt2  — light bands
BLUE    = RGBColor(0x1B, 0x5F, 0xA8)   # accent1 — primary / in progress
AMBER   = RGBColor(0xE0, 0x8A, 0x1E)   # accent2 — at risk / idle / tentative
ROSE    = RGBColor(0x9E, 0x3B, 0x4E)   # accent3 — late / blocked
GREEN   = RGBColor(0x4E, 0x7A, 0x3F)   # accent4 — done
SKY     = RGBColor(0x7F, 0xA9, 0xD4)   # accent6 — neutral highlight
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT_H = "Georgia"     # majorFont — only used inside the drawn timelines
FONT_B = "Segoe UI"    # minorFont

MONTHS_EN = {1: "January", 2: "February", 3: "March", 4: "April", 5: "May",
             6: "June", 7: "July", 8: "August", 9: "September", 10: "October",
             11: "November", 12: "December"}

_SEV_COLOUR = {"ok": GREEN, "warn": AMBER, "bad": ROSE, "none": MUTED}
_STATUS_COLOUR = {
    "Completed": GREEN, "Working on": BLUE, "Blocked": ROSE,
    "Not started": MUTED, "Cancelled": MUTED,
}
_ROWS_PER_SLIDE = 13     # header + 13 rows fits the 5.2" table placeholder


def _d(v):
    if not v:
        return None
    try:
        return datetime.date.fromisoformat(str(v)[:10])
    except Exception:
        return None


def _fmt(v) -> str:
    d = _d(v)
    return d.strftime("%d/%m/%Y") if d else "—"


def _long_date(d: datetime.date) -> str:
    return f"{d.day} {MONTHS_EN[d.month]} {d.year}"


# ── Template plumbing ─────────────────────────────────────────────────────────

_R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def _open() -> Presentation:
    prs = Presentation(str(TEMPLATE))
    ids = prs.slides._sldIdLst          # the .potx ships with zero slides,
    for sld in list(ids):               # but stay safe if that ever changes
        rid = sld.get(_R_ID)
        ids.remove(sld)
        if rid:
            try:
                prs.part.drop_rel(rid)
            except KeyError:
                pass
    return prs


def _layout(prs, name: str, fallback: str = "Tabula solita"):
    for m in prs.slide_masters:
        for l in m.slide_layouts:
            if l.name == name:
                return l
    for m in prs.slide_masters:
        for l in m.slide_layouts:
            if l.name == fallback:
                return l
    return prs.slide_masters[0].slide_layouts[0]


def _add(prs, layout_name: str):
    return prs.slides.add_slide(_layout(prs, layout_name))


def _ph(slide, idx: int):
    try:
        return slide.placeholders[idx]
    except KeyError:
        return None


def _set_ph(slide, idx: int, text, *, size=None, color=None, bold=None):
    """Fill a placeholder. No size/colour arguments → the LAYOUT formats it,
    which is the whole point of using the template's placeholders."""
    ph = _ph(slide, idx)
    if ph is None or not ph.has_text_frame:
        return None
    tf = ph.text_frame
    tf.text = str(text)
    if size or color or bold is not None:
        p = tf.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        if size:
            run.font.size = Pt(size)
        if color:
            run.font.color.rgb = color
        if bold is not None:
            run.font.bold = bold
    return ph


def _ph_paragraphs(slide, idx: int, blocks: list[dict]):
    """Rich multi-paragraph content inside one placeholder.

    blocks: [{"text", "size", "color", "bold", "space_after"}] — used for the
    per-person pages, where the two *Due contenuti* placeholders hold the four
    labelled lists as styled text instead of hand-drawn boxes.
    """
    ph = _ph(slide, idx)
    if ph is None or not ph.has_text_frame:
        return None
    tf = ph.text_frame
    tf.word_wrap = True
    first = True
    for b in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(b.get("space_after", 2))
        run = p.add_run()
        run.text = b["text"]
        run.font.size = Pt(b.get("size", 12))
        if b.get("color"):
            run.font.color.rgb = b["color"]
        if b.get("bold"):
            run.font.bold = True
    return ph


def _tidy(slide):
    """Remove placeholders left empty, so the deck has no 'Click to add' husks
    in edit view. What remains on a slide is exactly what carries content."""
    for shp in list(slide.shapes):
        if shp.is_placeholder and shp.has_text_frame and not shp.text_frame.text.strip():
            shp._element.getparent().remove(shp._element)


# ── Title slide (Tabula prima — always the first slide) ───────────────────────

def _title_slide(prs, title: str, subtitle: str, meta: str):
    s = _add(prs, "Tabula prima")
    _set_ph(s, 0, title)
    _set_ph(s, 1, subtitle)
    _set_ph(s, 10, meta)
    # idx 11 is the lab address block: inherit the layout's own text verbatim.
    lay = _layout(prs, "Tabula prima")
    lay_ph = next((p for p in lay.placeholders if p.placeholder_format.idx == 11), None)
    if lay_ph is not None and lay_ph.has_text_frame and _ph(s, 11) is not None:
        _ph(s, 11).text_frame.text = lay_ph.text_frame.text
    _tidy(s)
    return s


# ── Section divider (Sectio plena I–IV) ───────────────────────────────────────

_SECTIO = {1: "Sectio plena", 2: "Sectio plena II",
           3: "Sectio plena III", 4: "Sectio plena IV"}


def _divider_slide(prs, number: int, title: str, subtitle: str = ""):
    s = _add(prs, _SECTIO.get(number, "Sectio plena"))
    _set_ph(s, 12, f"{number:02d}")
    _set_ph(s, 0, title)
    if subtitle:
        _set_ph(s, 13, subtitle)
    _tidy(s)
    return s


# ── Tables (Tabella risultati — the native TABLE placeholder) ─────────────────

_TREE_COLS = [("ITEM", 6.1), ("STATUS", 1.5), ("DEADLINE", 1.3),
              ("OWNER / SUP.", 2.2), ("UPDATED", 1.2)]


def _cell(cell, text, *, size=10, bold=False, color=None, indent=0.0,
          align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.06 + indent)
    tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Emu(9525)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color


def _band(cell, colour):
    cell.fill.solid()
    cell.fill.fore_color.rgb = colour


def _tree_rows(prs, title, subtitle, rows, footer=""):
    """One 'Tabella risultati' slide per page of rows.

    Title placeholder = title; the small caption placeholder under the table
    carries the subtitle; the layout's FOOTER placeholder carries the footnote.
    Hierarchy inside the table: indentation + a light band on group rows.
    """
    slides = []
    pages = [rows[i:i + _ROWS_PER_SLIDE] for i in range(0, len(rows), _ROWS_PER_SLIDE)] or [[]]

    for pi, page in enumerate(pages):
        s = _add(prs, "Tabella risultati")
        _set_ph(s, 0, title if pi == 0 else f"{title} (cont.)")
        # python-pptx does not clone FOOTER placeholders onto new slides, so the
        # footnote rides in the caption placeholder together with the subtitle.
        caption = f"{subtitle}   ·   {footer}" if footer else subtitle
        _set_ph(s, 13, caption, size=10, color=MUTED)

        if not page:
            tbl_ph = _ph(s, 12)
            if tbl_ph is not None:
                gf = tbl_ph.insert_table(rows=1, cols=1)
                _cell(gf.table.cell(0, 0), "Nothing to show.", size=11, color=MUTED)
            _tidy(s)
            slides.append(s)
            continue

        tbl_ph = _ph(s, 12)
        gf = tbl_ph.insert_table(rows=len(page) + 1, cols=len(_TREE_COLS))
        table = gf.table
        table.first_row = True
        for i, (_, w) in enumerate(_TREE_COLS):
            table.columns[i].width = Inches(w)

        for i, (label, _) in enumerate(_TREE_COLS):
            _cell(table.cell(0, i), label, size=9, bold=True)
        table.rows[0].height = Inches(0.28)

        for ri, r in enumerate(page, start=1):
            table.rows[ri].height = Inches(0.28)
            kind = r.get("kind", "task")
            lvl = int(r.get("level", 0))

            if kind in ("deliverable", "group"):
                for ci in range(len(_TREE_COLS)):
                    _band(table.cell(ri, ci), RULE)
                marker, size, bold, colour = "▸", 11, True, BLUE
            elif kind == "subtask":
                marker, size, bold, colour = "›", 9, False, MUTED
            else:
                marker, size, bold, colour = "•", 10, False, INK

            _cell(table.cell(ri, 0), f"{marker} {r.get('name', '')}",
                  size=size, bold=bold, color=colour, indent=0.22 * lvl)
            status = r.get("status") or ""
            _cell(table.cell(ri, 1), status, size=8, bold=bool(status),
                  color=_STATUS_COLOUR.get(status, MUTED))
            _cell(table.cell(ri, 2), _fmt(r["deadline"]) if r.get("deadline") else "",
                  size=8, color=INK)
            _cell(table.cell(ri, 3), r.get("people") or "", size=8, color=MUTED)
            _cell(table.cell(ri, 4), r.get("fresh") or "", size=8,
                  bold=r.get("sev") in ("warn", "bad"),
                  color=_SEV_COLOUR.get(r.get("sev", "none"), MUTED))

        _tidy(s)
        slides.append(s)
    return slides


_ENGAGE_COLS = [("PERSON", 3.2), ("ACTIVITY", 4.0), ("UPDATES", 1.2),
                ("COMMENTS", 1.3), ("DAYS SEEN", 1.2), ("LAST SEEN", 1.4)]
_ENGAGE_COLS_NO_LOGIN = [("PERSON", 3.6), ("ACTIVITY", 5.2), ("UPDATES", 1.6),
                         ("COMMENTS", 1.9)]
_ENGAGE_STATE = {"active": BLUE, "quiet": AMBER, "absent": MUTED, "silent": MUTED}


def _engagement_slide(prs, engagement, period):
    """Who is using the board — one bar per person, over the deck's period.

    Deliberately NOT a leaderboard: no rank numbers, no red. The bar is the
    whole message — active people carry a visible one, the quiet rows simply
    fade to a lowercase note. Enough to be read across a room in two seconds,
    not enough to pillory anyone whose task this month is a three-week
    simulation run.
    """
    people = engagement.get("people") or []
    if not people:
        return None
    tot = engagement.get("totals", {})
    has_logins = engagement.get("logins_available", False)
    cols = _ENGAGE_COLS if has_logins else _ENGAGE_COLS_NO_LOGIN
    top = max(tot.get("top", 0), 1)

    caption = (f"{period} · {tot.get('active', 0)} of {tot.get('people', 0)} people "
               f"contributed · {tot.get('updates', 0)} updates, "
               f"{tot.get('comments', 0)} comments")
    footer = ("Activity = status changes + comments — the acts that leave a trace "
              "others can read. Signing in is not scored."
              if has_logins else
              "Activity = status changes + comments. Sign-in tracking is not "
              "enabled yet, so attendance is not shown.")

    slides = []
    pages = [people[i:i + _ROWS_PER_SLIDE]
             for i in range(0, len(people), _ROWS_PER_SLIDE)] or [[]]
    for pi, page in enumerate(pages):
        s = _add(prs, "Tabella risultati")
        _set_ph(s, 0, "Who is using the board" if pi == 0
                else "Who is using the board (cont.)")
        _set_ph(s, 13, f"{caption}   ·   {footer}", size=10, color=MUTED)

        gf = _ph(s, 12).insert_table(rows=len(page) + 1, cols=len(cols))
        table = gf.table
        table.first_row = True
        for i, (label, w) in enumerate(cols):
            table.columns[i].width = Inches(w)
            _cell(table.cell(0, i), label, size=9, bold=True,
                  align=PP_ALIGN.CENTER if i >= 2 else PP_ALIGN.LEFT)
        table.rows[0].height = Inches(0.28)

        for ri, p in enumerate(page, start=1):
            table.rows[ri].height = Inches(0.28)
            state = p.get("state", "absent")
            colour = _ENGAGE_STATE.get(state, MUTED)
            n = p.get("contributions", 0)

            _cell(table.cell(ri, 0), p.get("name", "—"), size=10,
                  bold=state == "active", color=INK if state == "active" else MUTED)

            if n:
                # 12 blocks at the top of the lab; at least one when non-zero,
                # so a single update still shows as a mark rather than nothing.
                bar = "■" * max(1, round(12 * n / top))
                _cell(table.cell(ri, 1), f"{bar}  {n}", size=10, bold=True, color=colour)
            else:
                # Never claim someone "signs in" when sign-ins are not tracked:
                # without login_events the only defensible statement is silence.
                if not has_logins:
                    label = "no updates"
                elif state == "quiet":
                    label = "signs in, no updates"
                else:
                    label = "not seen"
                _cell(table.cell(ri, 1), label, size=9, color=colour)

            _cell(table.cell(ri, 2), p.get("updates", 0) or "—", size=9,
                  color=INK if p.get("updates") else MUTED, align=PP_ALIGN.CENTER)
            _cell(table.cell(ri, 3), p.get("comments", 0) or "—", size=9,
                  color=INK if p.get("comments") else MUTED, align=PP_ALIGN.CENTER)
            if has_logins:
                days_in = p.get("days_in") or 0
                _cell(table.cell(ri, 4), days_in or "—", size=9,
                      color=INK if days_in else MUTED, align=PP_ALIGN.CENTER)
                last = p.get("last_seen")
                _cell(table.cell(ri, 5),
                      last.strftime("%d/%m/%Y") if last else "—",
                      size=9, color=MUTED, align=PP_ALIGN.CENTER)
        _tidy(s)
        slides.append(s)
    return slides


def _stats_slide(prs, title, subtitle, tiles, footer=""):
    """The at-a-glance numbers as a small native table: value row + label row.
    Replaces the old hand-drawn coloured tiles."""
    s = _add(prs, "Tabella risultati")
    _set_ph(s, 0, title)
    caption = f"{subtitle}   ·   {footer}" if footer else subtitle
    _set_ph(s, 13, caption, size=10, color=MUTED)

    tbl_ph = _ph(s, 12)
    gf = tbl_ph.insert_table(rows=2, cols=max(len(tiles), 1))
    table = gf.table
    table.first_row = False
    table.horz_banding = False
    for i, (label, value, colour) in enumerate(tiles):
        table.rows[0].height = Inches(0.9)
        table.rows[1].height = Inches(0.4)
        _cell(table.cell(0, i), str(value), size=40, bold=True, color=colour,
              align=PP_ALIGN.CENTER)
        _cell(table.cell(1, i), label.upper(), size=10, bold=True, color=MUTED,
              align=PP_ALIGN.CENTER)
        _band(table.cell(0, i), WHITE)
        _band(table.cell(1, i), WHITE)
    _tidy(s)
    return s


# ── Per-person page (Due contenuti) ───────────────────────────────────────────

def _list_block(label, colour, items, sub):
    blocks = [{"text": f"{label} · {len(items)}", "size": 12, "bold": True,
               "color": colour, "space_after": 3}]
    for it in items[:6]:
        blocks.append({"text": f"• {it.get('name', '')}", "size": 11, "color": INK})
        extra = sub(it)
        if extra:
            blocks.append({"text": f"   {extra}", "size": 8.5, "color": MUTED,
                           "space_after": 3})
    if len(items) > 6:
        blocks.append({"text": f"… and {len(items) - 6} more", "size": 9, "color": MUTED})
    if not items:
        blocks.append({"text": "—", "size": 10, "color": MUTED})
    blocks.append({"text": "", "size": 6, "space_after": 6})
    return blocks


def _person_slide(prs, person, period):
    s = _add(prs, "Due contenuti")
    _set_ph(s, 0, person.get("name", "—"))
    dom = person.get("dominant_project") or ""
    ctx = f"{period} · {person.get('active', 0)} active tasks"
    if dom:
        ctx += f" · mainly {dom}"

    left = ([{"text": ctx, "size": 10, "color": MUTED, "space_after": 8}]
            + (_list_block("COMPLETED", GREEN, person.get("completed", []),
                        lambda it: it.get("_project", ""))
            + _list_block("MOVED", BLUE, person.get("moved", []),
                          lambda it: f"{it.get('_project', '')} · {it.get('status', '')}")))
    right = (_list_block("BLOCKED", ROSE, person.get("blocked", []),
                         lambda it: f"{it.get('_project', '')} · idle {it.get('_stale', '?')}d")
             + _list_block("IDLE", AMBER, person.get("stale", []),
                           lambda it: f"{it.get('_project', '')} · idle {it.get('_stale', '?')}d"))
    _ph_paragraphs(s, 12, left)
    _ph_paragraphs(s, 13, right)
    _tidy(s)
    return s


# ── Timelines (the one drawn exception — data figures, theme colours only) ────

def _text(slide, x, y, w, h, text, *, size=10, bold=False, color=INK,
          align=PP_ALIGN.LEFT, wrap=False):
    box = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT_B
    return box


def _rect(slide, x, y, w, h, fill):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
                                Emu(int(w)), Emu(int(h)))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _assign_lanes(xs, label_width, gap=None, max_lanes=4):
    """Proximity-aware label lanes for timeline markers.

    ``xs`` are marker centre x positions (EMU). Lanes are 0=above, 1=below,
    2=above-outer, 3=below-outer. Each marker prefers the alternating
    above/below rhythm, but a lane is only granted if the previous label on
    that lane leaves horizontal room — near-coincident dates escalate to the
    outer lanes instead of overlapping."""
    if gap is None:
        gap = Inches(0.08)
    half = label_width / 2
    right_edge = [None] * max_lanes
    lanes = [0] * len(xs)
    for k, i in enumerate(sorted(range(len(xs)), key=lambda j: xs[j])):
        prefs = [0, 1, 2, 3] if k % 2 == 0 else [1, 0, 3, 2]
        prefs = [p for p in prefs if p < max_lanes]
        for lane in prefs:
            if right_edge[lane] is None or xs[i] - half >= right_edge[lane] + gap:
                break
        else:  # every lane crowded: take the one that frees up first
            lane = min(prefs, key=lambda p: right_edge[p])
        right_edge[lane] = xs[i] + half
        lanes[i] = lane
    return lanes


def _timeline_slide(prs, confs, months):
    """Conference submissions on a horizontal axis (Solo titolo + drawing)."""
    s = _add(prs, "Solo titolo")
    _set_ph(s, 0, "Submission timeline")
    W = prs.slide_width
    _text(s, Inches(0.5), prs.slide_height - Inches(0.45), W - Inches(1.0), Inches(0.3),
          f"Next {months} months · red ≤30d · amber ≤90d", size=9, color=MUTED)
    x0, x1 = Inches(1.0), W - Inches(1.0)
    y = Inches(3.6)
    span = max(months * 30.44, 1)
    today = datetime.date.today()

    _rect(s, x0, y, x1 - x0, Emu(19050), RULE)
    for m in range(months + 1):
        d = today + datetime.timedelta(days=int(30.44 * m))
        fx = x0 + (x1 - x0) * (30.44 * m) / span
        _rect(s, fx, y - Inches(0.07), Emu(9525), Inches(0.14), RULE)
        if m % max(1, months // 6) == 0:
            _text(s, fx - Inches(0.35), y + Inches(0.16), Inches(0.7), Inches(0.2),
                  d.strftime("%b %y"), size=8, color=MUTED, align=PP_ALIGN.CENTER)

    xs = [int(x0 + (x1 - x0) * min(max(c["_days"], 0) / span, 1.0)) for c in confs]
    lanes = _assign_lanes(xs, Inches(2.0))
    lane_y = {0: y - Inches(0.95), 1: y + Inches(0.5),
              2: y - Inches(1.55), 3: y + Inches(1.1)}
    for c, fx, lane in zip(confs, xs, lanes):
        colour = ROSE if c["_days"] <= 30 else (AMBER if c["_days"] <= 90 else BLUE)
        _rect(s, fx - Emu(28575), y - Inches(0.12), Emu(57150), Inches(0.28), colour)
        if lane == 2:   # leader tick up to the outer lane
            _rect(s, fx - Emu(4763), y - Inches(1.06), Emu(9525), Inches(0.94), RULE)
        elif lane == 3:
            _rect(s, fx - Emu(4763), y + Inches(0.16), Emu(9525), Inches(0.94), RULE)
        ly = lane_y[lane]
        _text(s, fx - Inches(0.95), ly, Inches(1.9), Inches(0.24),
              c["_label"], size=10, bold=True, color=colour, align=PP_ALIGN.CENTER)
        _text(s, fx - Inches(0.95), ly + Inches(0.22), Inches(1.9), Inches(0.2),
              f"{_fmt(c.get('submission_deadline'))} · {len(c['_papers'])} draft"
              f"{'' if len(c['_papers']) == 1 else 's'}",
              size=8, color=MUTED, align=PP_ALIGN.CENTER)
    _tidy(s)
    return s


def _publication_timeline(prs, pack):
    """Published (left of today) vs drafts and targets (right) — the pipeline."""
    s = _add(prs, "Solo titolo")
    _set_ph(s, 0, "Publication pipeline")
    W = prs.slide_width
    _text(s, Inches(0.5), prs.slide_height - Inches(0.45), W - Inches(1.0), Inches(0.3),
          "Green: published (Scopus) · blue: approved targets · amber: tentative topics",
          size=9, color=MUTED)
    x0, x1 = Inches(1.0), W - Inches(1.0)
    y = Inches(3.6)
    today = datetime.date.today()

    pubs = pack.get("publications") or {}
    by_year: dict[int, int] = {}
    for kind in ("journals_by_year", "conferences_by_year"):
        for yr, items in (pubs.get(kind) or {}).items():
            by_year[int(yr)] = by_year.get(int(yr), 0) + len(items)

    targets = []
    conf_by_key = {}
    for cconf in pack.get("conferences", []):
        lbl = (cconf.get("acronym") or cconf.get("name") or "").strip()
        key = f"{lbl} {cconf['year']}".strip() if cconf.get("year") else lbl
        conf_by_key[key.lower()] = cconf
    for p in pack.get("conf_papers", []):
        cconf = conf_by_key.get((p.get("_target") or "").lower())
        sub = _d((cconf or {}).get("submission_deadline")) or _d(p.get("deadline"))
        if sub and sub >= today:
            targets.append((sub, p.get("_target") or p.get("name", ""),
                            bool(p.get("_tentative"))))
    for dr in pack.get("paper_drafts", []):
        dl = _d(dr.get("deadline"))
        if dl and dl >= today:
            targets.append((dl, dr.get("name", ""), False))
    targets.sort()

    years_back = sorted(by_year)[-4:] if by_year else []
    span_start = (datetime.date(years_back[0], 1, 1) if years_back
                  else today - datetime.timedelta(days=365))
    span_end = max([t[0] for t in targets], default=today) + datetime.timedelta(days=90)
    past_days = max((today - span_start).days, 0)
    future_days = max((span_end - today).days, 1)
    # years of history would crush every future target into the right edge:
    # give the future (the actionable part) at least half of the axis
    past_frac = min(past_days / max(past_days + future_days, 1), 0.5)
    xm = x0 + (x1 - x0) * past_frac

    def fx(d):
        if d <= today:
            if past_days == 0:
                return x0
            return x0 + (xm - x0) * min(max((d - span_start).days / past_days, 0.0), 1.0)
        return xm + (x1 - xm) * min((d - today).days / future_days, 1.0)

    _rect(s, x0, y, x1 - x0, Emu(19050), RULE)
    tx = fx(today)
    # the bar stops short of both label lanes; the caption rides beside it
    _rect(s, tx, y - Inches(0.36), Emu(19050), Inches(0.72), INK)
    _text(s, tx + Inches(0.05), y - Inches(0.35), Inches(0.6), Inches(0.2),
          "today", size=8, bold=True, color=INK)

    for yr in years_back:
        n = by_year[yr]
        cx = fx(datetime.date(yr, 7, 1))
        h = Inches(0.16) * min(n, 6)
        _rect(s, cx - Inches(0.14), y - h - Emu(19050), Inches(0.28), h, GREEN)
        _text(s, cx - Inches(0.4), y + Inches(0.12), Inches(0.8), Inches(0.22),
              str(yr), size=9, color=MUTED, align=PP_ALIGN.CENTER)
        _text(s, cx - Inches(0.4), y - h - Inches(0.32), Inches(0.8), Inches(0.22),
              str(n), size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # same submission date (several drafts for one venue) → one marker, "+N"
    grouped = []
    for sub, label, tentative in targets:
        if grouped and grouped[-1][0] == sub:
            grouped[-1][1].append(label)
            grouped[-1][2] = grouped[-1][2] and tentative
        else:
            grouped.append([sub, [label], tentative])
    shown = [(sub, names[0] if len(names) == 1
              else f"{names[0][:18]} +{len(names) - 1}", tent)
             for sub, names, tent in grouped[:8]]
    txs = [int(fx(sub)) for sub, _, _ in shown]
    lanes = _assign_lanes(txs, Inches(2.1))
    lane_y = {0: y - Inches(0.95), 1: y + Inches(0.42),
              2: y - Inches(1.5), 3: y + Inches(1.0)}
    for (sub, label, tentative), cx, lane in zip(shown, txs, lanes):
        colour = AMBER if tentative else BLUE
        _rect(s, cx - Emu(28575), y - Inches(0.12), Emu(57150), Inches(0.28), colour)
        if lane == 2:   # leader tick up to the outer lane
            _rect(s, cx - Emu(4763), y - Inches(1.03), Emu(9525), Inches(0.91), RULE)
        elif lane == 3:
            _rect(s, cx - Emu(4763), y + Inches(0.16), Emu(9525), Inches(0.84), RULE)
        ly = lane_y[lane]
        _text(s, cx - Inches(1.0), ly, Inches(2.0), Inches(0.22),
              label[:26], size=9, bold=True, color=colour, align=PP_ALIGN.CENTER)
        _text(s, cx - Inches(1.0), ly + Inches(0.2), Inches(2.0), Inches(0.2),
              _fmt(sub.isoformat()) + (" · tentative" if tentative else ""),
              size=8, color=MUTED, align=PP_ALIGN.CENTER)

    if not by_year and not targets:
        _text(s, Inches(1.0), Inches(2.4), W - Inches(2.0), Inches(0.4),
              "No publications recorded and no targets yet — pick a conference "
              "under Conference Paper Drafts to start the pipeline.",
              size=13, color=MUTED, wrap=True)
    _tidy(s)
    return s


# ── Deck 1: lab meeting ───────────────────────────────────────────────────────

def build_meeting_deck(pack: dict) -> BytesIO:
    prs = _open()
    since, until = pack["since"], pack["until"]
    period = f"{_long_date(since)} — {_long_date(until)}"

    _title_slide(prs, "Lab meeting", "Status review", period)

    # 1 · Deliverables due soon
    horizon = pack.get("horizon_months", 3)
    ups = pack.get("upcoming", [])
    _divider_slide(prs, 1, "Deliverables due soon",
                   f"Next {horizon} months · {len(ups)} deliverables")
    rows = []
    for d in ups:
        days = d["_days"]
        sev = "bad" if days < 0 else ("warn" if days <= 30 else "ok")
        when = (f"{abs(days)}d overdue" if days < 0
                else ("today" if days == 0 else f"in {days}d"))
        rows.append({"level": 0, "kind": "deliverable",
                     "name": f"[{d['_project']}] {d.get('name', '')}",
                     "status": d.get("status") or "Not started",
                     "deadline": d.get("deadline"), "people": d["_people"],
                     "fresh": when, "sev": sev})
        rows.append({"level": 1, "kind": "task",
                     "name": f"{d['_done']}/{d['_total']} tasks completed · {d['_pct']}%",
                     "status": "", "deadline": None, "people": "",
                     "fresh": "", "sev": "none"})
    _tree_rows(prs, "Deliverables due soon",
               f"Deadline within {horizon} months · sorted by urgency", rows,
               footer="The UPDATED column here is time to deadline, not freshness.")

    # 2 · Projects
    trees = pack.get("trees", [])
    _divider_slide(prs, 2, "Projects",
                   f"{len(trees)} active projects · deliverables, tasks and subtasks")
    for t in trees:
        head = f"{t['acronym']} — {t['name']}" if t["acronym"] else t["name"]
        c = t["counts"]
        _tree_rows(prs, head,
                   f"{c['deliverables']} deliverables · {c['tasks']} tasks",
                   t["rows"],
                   footer="UPDATED = days since anyone last touched the item.")

    # 3 · Individual work
    delta = pack.get("delta", {})
    people = delta.get("by_person", [])
    tot = delta.get("totals", {})
    _divider_slide(prs, 3, "Individual work", f"{period} · {len(people)} people")
    _stats_slide(prs, "What changed", period,
                 [("completed", tot.get("completed", 0), GREEN),
                  ("moved", tot.get("moved", 0), BLUE),
                  ("blocked", tot.get("blocked", 0), ROSE),
                  ("idle", tot.get("stale", 0), AMBER)],
                 footer="Idle = no update within the staleness threshold.")

    engagement = pack.get("engagement")
    if engagement:
        _engagement_slide(prs, engagement, pack.get("engagement_period") or period)

    for person in people:
        _person_slide(prs, person, period)

    # 4 · Next deadlines
    confs = pack.get("conferences", [])
    ch = pack.get("conf_horizon_months", 12)
    _divider_slide(prs, 4, "Next deadlines",
                   f"Conference submissions · next {ch} months")
    if confs:
        _timeline_slide(prs, confs, ch)
        crows = []
        for c in confs:
            days = c["_days"]
            sev = "bad" if days <= 30 else ("warn" if days <= 90 else "ok")
            loc = f" · {c.get('location')}" if c.get("location") else ""
            crows.append({"level": 0, "kind": "deliverable",
                          "name": f"{c['_label']}{loc}",
                          "status": "", "deadline": c.get("submission_deadline"),
                          "people": (c.get("rank") or ""),
                          "fresh": f"in {days}d" if days >= 0 else f"{abs(days)}d ago",
                          "sev": sev})
            if c["_papers"]:
                crows.extend(c["_papers"])
            else:
                crows.append({"level": 1, "kind": "subtask",
                              "name": "no paper draft yet", "status": "",
                              "deadline": None, "people": "", "fresh": "",
                              "sev": "none"})
        _tree_rows(prs, "Conferences and paper drafts",
                   "Submission deadline, and the papers targeting each", crows,
                   footer="A conference with no draft is a decision waiting to be made.")
    else:
        _tree_rows(prs, "Next deadlines",
                   "No conference with a submission deadline in the window", [])

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Deck 2: project review (cumulative) ───────────────────────────────────────

def build_review_deck(review: dict) -> BytesIO:
    prs = _open()
    proj = review.get("project", {})
    name = proj.get("name", "Project")
    acr = proj.get("acronym") or proj.get("identifier") or ""
    today = datetime.date.today()
    period = review.get("period_label") or f"Status as of {_long_date(today)}"

    _title_slide(prs, acr or name, name if acr else "Progress status", period)

    t = review.get("totals", {})
    done, total = t.get("completed", 0), t.get("total", 0)
    pct = round(100 * done / total) if total else 0
    _stats_slide(prs, "Progress status", f"{name} · {period}",
                 [("completion", f"{pct}%", GREEN if pct >= 66 else BLUE),
                  ("total tasks", total, INK),
                  ("overdue", t.get("overdue", 0), ROSE),
                  ("at risk", t.get("at_risk", 0), AMBER)],
                 footer=f"CUP {proj.get('cup') or '—'} · {proj.get('funding_agency') or ''}")

    for d in review.get("deliverables", []):
        dt = d.get("totals", {})
        rows = []
        for label, colour_kind, items, sub in (
            ("COMPLETED", "group", d.get("completed", []),
             lambda it: f"closed {_fmt(it.get('completion_date'))}"),
            ("IN PROGRESS", "group", d.get("in_progress", []),
             lambda it: f"due {_fmt(it.get('deadline'))}"),
            ("RISKS", "group", d.get("risks", []),
             lambda it: it.get("_why", "")),
        ):
            rows.append({"level": 0, "kind": "group", "name": f"{label} · {len(items)}",
                         "status": "", "deadline": None, "people": "",
                         "fresh": "", "sev": "none"})
            for it in items[:7]:
                rows.append({"level": 1, "kind": "task", "name": it.get("name", ""),
                             "status": it.get("status") or "",
                             "deadline": it.get("deadline"),
                             "people": sub(it), "fresh": "", "sev": "none"})
        _tree_rows(prs, d.get("name", "Deliverable"),
                   f"{d.get('type', '')} · due {_fmt(d.get('deadline'))} · "
                   f"{dt.get('completed', 0)}/{dt.get('total', 0)} tasks · "
                   f"{dt.get('pct', 0)}% complete",
                   rows,
                   footer=f"{name} · generated {_long_date(today)}")

    orphans = review.get("no_deliverable", [])
    if orphans:
        rows = [{"level": 0, "kind": "task", "name": it.get("name", ""),
                 "status": it.get("status") or "",
                 "deadline": it.get("deadline"),
                 "people": "", "fresh": "", "sev": "none"} for it in orphans]
        _tree_rows(prs, "Tasks not linked to a deliverable",
                   f"{len(orphans)} tasks", rows)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Deck 3: personal status ───────────────────────────────────────────────────

def build_my_status_deck(pack: dict) -> BytesIO:
    prs = _open()
    user = pack.get("user", {})
    name = user.get("name") or user.get("email") or "—"
    since, until = pack["since"], pack["until"]
    period = f"{_long_date(since)} — {_long_date(until)}"
    c = pack.get("counts", {})

    _title_slide(prs, name, "Status review", period)

    pubs = pack.get("publications") or {}
    pub_total = (pubs.get("totals") or {}).get("all", 0)
    _stats_slide(prs, "At a glance", period,
                 [("completed", c.get("completed", 0), GREEN),
                  ("active tasks", c.get("active", 0), BLUE),
                  ("papers in pipeline", c.get("papers", 0), INK),
                  ("publications", pub_total, GREEN if pub_total else AMBER)],
                 footer="Papers in pipeline = journal drafts + conference "
                        "targets (tentative included).")

    _tree_rows(prs, "Completed in the period", period, pack.get("completed", []),
               footer="Completion dates from the task records.")

    _publication_timeline(prs, pack)

    rows = []
    for d in pack.get("paper_drafts", []):
        rows.append({"level": 0, "kind": "deliverable",
                     "name": f"[{d.get('_project', '')}] {d.get('name', '')}",
                     "status": d.get("status") or "Not started",
                     "deadline": d.get("deadline"),
                     "people": "journal / deliverable", "fresh": "", "sev": "none"})
    for p in pack.get("conf_papers", []):
        rows.append({"level": 0, "kind": "deliverable",
                     "name": p.get("name") or "—",
                     "status": p.get("status") or "Not started",
                     "deadline": p.get("deadline"),
                     "people": p.get("_target") or "no conference yet",
                     "fresh": "TENT." if p.get("_tentative") else "",
                     "sev": "warn" if p.get("_tentative") else "none"})
    _tree_rows(prs, "Papers", "Journal drafts and conference targets", rows,
               footer="TENT. = tentative topic, not yet approved by the supervisor.")

    _tree_rows(prs, "Active tasks",
               f"{c.get('active', 0)} open items by project",
               pack.get("task_rows", []),
               footer="UPDATED = days since anyone last touched the item.")

    if user.get("is_phd_student"):
        _phd_slide(prs, user)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _phd_slide(prs, user):
    """The PhD clock as a key message (Messaggio chiave layout) — no drawing."""
    s = _add(prs, "Messaggio chiave")
    start, end = _d(user.get("phd_start_date")), _d(user.get("phd_end_date"))
    if not start or not end or end <= start:
        _ph_paragraphs(s, 10, [{"text": "PhD start/end dates are not configured "
                                        "(Admin Panel → Users).",
                                "size": 16, "color": MUTED}])
        _tidy(s)
        return s
    today = datetime.date.today()
    total = (end - start).days
    pct = min(max((today - start).days / total, 0.0), 1.0)
    years = max(1, round(total / 365.25))
    year_now = min(years, max(1, int((today - start).days // 365.25) + 1))
    _ph_paragraphs(s, 10, [
        {"text": f"PhD: year {year_now} of {years}", "size": 32, "bold": True,
         "color": BLUE, "space_after": 8},
        {"text": f"{_fmt(start.isoformat())} → {_fmt(end.isoformat())} · "
                 f"{pct * 100:.0f}% of the time elapsed",
         "size": 14, "color": MUTED},
    ])
    _tidy(s)
    return s
